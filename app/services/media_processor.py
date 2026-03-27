from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from app.core.config import settings

def process_pptx(file_path: str, filename: str) -> list[dict]:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
                
    if not text.strip():
        raise ValueError(f"File slide '{filename}' hoàn toàn rỗng hoặc chỉ chứa ảnh.")

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
    )
    chunks = text_splitter.split_text(text)
    return [{"page_content": chunk, "metadata": {"filename": filename}} for chunk in chunks]
